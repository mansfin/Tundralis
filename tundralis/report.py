"""PowerPoint report generation using python-pptx."""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)

# ─── Brand colors ─────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)
MID_GRAY = RGBColor(0x8C, 0x9B, 0xB2)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)

# ─── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class ReportBuilder:
    """Builds the full KDA PowerPoint report."""

    def __init__(self, results, narratives, charts: dict[str, bytes]):
        self.r = results
        self.n = narratives
        self.charts = charts
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    # ── Layout helpers ────────────────────────────────────────────────────────

    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]  # blank
        return self.prs.slides.add_slide(layout)

    def _add_rect(self, slide, left, top, width, height, fill_color: RGBColor, alpha: Optional[int] = None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height,
        )
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color
        return shape

    def _add_textbox(
        self,
        slide,
        text: str,
        left, top, width, height,
        font_size: int = 12,
        bold: bool = False,
        color: RGBColor = DARK_BLUE,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        word_wrap: bool = True,
        italic: bool = False,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return txBox

    def _add_image_bytes(self, slide, img_bytes: bytes, left, top, width, height=None):
        img_stream = io.BytesIO(img_bytes)
        if height:
            slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(img_stream, left, top, width=width)

    def _header_bar(self, slide, title: str, subtitle: str = ""):
        """Dark blue header bar at top of content slides."""
        self._add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
        self._add_textbox(
            slide, title,
            Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
            font_size=22, bold=True, color=WHITE,
        )
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                Inches(0.4), Inches(0.65), Inches(10), Inches(0.35),
                font_size=11, color=TEAL,
            )
        # Teal accent line
        self._add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.04), TEAL)
        # Footer
        self._add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), DARK_BLUE)
        self._add_textbox(
            slide, "tundralis.com  ·  Confidential",
            Inches(0.3), Inches(7.22), Inches(5), Inches(0.25),
            font_size=7, color=MID_GRAY,
        )
        self._add_textbox(
            slide, "Key Driver Analysis",
            Inches(9), Inches(7.22), Inches(4), Inches(0.25),
            font_size=7, color=MID_GRAY, align=PP_ALIGN.RIGHT,
        )

    # ── Slides ────────────────────────────────────────────────────────────────

    def _slide_title(self):
        slide = self._blank_slide()

        # Full background
        self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)

        # Teal accent band (left side)
        self._add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, TEAL)

        # Teal bottom strip
        self._add_rect(slide, 0, Inches(6.8), SLIDE_W, Inches(0.7), TEAL)

        # Brand name (top left)
        self._add_textbox(
            slide, "TUNDRALIS",
            Inches(0.4), Inches(0.3), Inches(4), Inches(0.5),
            font_size=14, bold=True, color=TEAL,
        )
        self._add_textbox(
            slide, "tundralis.com",
            Inches(0.4), Inches(0.75), Inches(4), Inches(0.35),
            font_size=10, color=MID_GRAY,
        )

        # Main title
        self._add_textbox(
            slide, "Key Driver Analysis",
            Inches(0.4), Inches(1.8), Inches(12), Inches(1.2),
            font_size=40, bold=True, color=WHITE,
        )

        # Subtitle
        target_label = self.r.target.replace("_", " ").title()
        self._add_textbox(
            slide, f"Drivers of {target_label}",
            Inches(0.4), Inches(3.0), Inches(12), Inches(0.7),
            font_size=22, color=TEAL,
        )

        # Meta info
        meta_lines = (
            f"Sample size: {self.r.n_obs:,} respondents  ·  "
            f"Predictors analyzed: {len(self.r.predictors)}  ·  "
            f"Model R² = {self.r.meta['r_squared']:.3f}"
        )
        self._add_textbox(
            slide, meta_lines,
            Inches(0.4), Inches(3.9), Inches(12), Inches(0.4),
            font_size=11, color=MID_GRAY,
        )

        # Bottom strip text
        self._add_textbox(
            slide, "CONFIDENTIAL  ·  FOR CLIENT USE ONLY",
            Inches(0.4), Inches(6.87), Inches(12), Inches(0.3),
            font_size=9, bold=True, color=DARK_BLUE,
        )

    def _slide_exec_summary(self, summary_text: str):
        slide = self._blank_slide()
        self._header_bar(slide, "Executive Summary")

        # Summary box
        self._add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(4.2), LIGHT_GRAY)
        self._add_textbox(
            slide, summary_text,
            Inches(0.7), Inches(1.5), Inches(12.0), Inches(3.8),
            font_size=14, color=DARK_BLUE, word_wrap=True,
        )

        # Key stats row
        stats = [
            ("Sample Size", f"{self.r.n_obs:,}"),
            ("Predictors", str(len(self.r.predictors))),
            ("Model R²", f"{self.r.meta['r_squared']:.3f}"),
            ("Adj. R²", f"{self.r.meta['adj_r_squared']:.3f}"),
            ("Sig. Drivers", str(self.r.regression.coefficients["significant"].sum())),
        ]
        box_w = Inches(2.3)
        start_x = Inches(0.4)
        for i, (label, value) in enumerate(stats):
            x = start_x + i * box_w
            self._add_rect(slide, x, Inches(5.7), Inches(2.1), Inches(0.95), DARK_BLUE)
            self._add_textbox(
                slide, value,
                x + Inches(0.1), Inches(5.75), Inches(1.9), Inches(0.5),
                font_size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
            )
            self._add_textbox(
                slide, label,
                x + Inches(0.1), Inches(6.25), Inches(1.9), Inches(0.3),
                font_size=9, color=MID_GRAY, align=PP_ALIGN.CENTER,
            )

    def _slide_methodology(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Methodology", "How we identified what matters most")

        steps = [
            ("01", "Data Collection", "Survey responses collected across key experience dimensions."),
            ("02", "Correlation Analysis", "Pearson and Spearman correlations identify linear and monotonic relationships between each driver and the outcome."),
            ("03", "OLS Regression", "Ordinary Least Squares regression with standardized coefficients (β) quantifies each driver's unique contribution."),
            ("04", "Relative Importance", "Shapley value decomposition allocates the model's R² equitably across predictors, accounting for multicollinearity."),
            ("05", "Priority Matrix", "Drivers are mapped onto an Importance × Performance quadrant to identify quick wins and areas of risk."),
        ]

        for i, (num, title, desc) in enumerate(steps):
            y = Inches(1.35) + i * Inches(1.08)
            # Number bubble
            self._add_rect(slide, Inches(0.35), y, Inches(0.55), Inches(0.55), TEAL)
            self._add_textbox(
                slide, num,
                Inches(0.35), y, Inches(0.55), Inches(0.55),
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            )
            # Title + description
            self._add_textbox(
                slide, title,
                Inches(1.05), y, Inches(5), Inches(0.3),
                font_size=11, bold=True, color=DARK_BLUE,
            )
            self._add_textbox(
                slide, desc,
                Inches(1.05), y + Inches(0.28), Inches(11.8), Inches(0.55),
                font_size=9.5, color=MID_GRAY, word_wrap=True,
            )

    def _slide_importance_chart(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Key Drivers Ranked by Importance", "Shapley relative importance — share of explained variance")

        if "importance_bar" in self.charts:
            self._add_image_bytes(
                slide, self.charts["importance_bar"],
                Inches(0.4), Inches(1.25),
                width=Inches(12.5),
            )

    def _slide_quadrant(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Priority Matrix", "Importance vs. Performance — where to focus")

        if "quadrant" in self.charts:
            self._add_image_bytes(
                slide, self.charts["quadrant"],
                Inches(0.5), Inches(1.2),
                width=Inches(8.5),
            )

        # Legend / interpretation
        quadrant_notes = [
            ("Priority Fixes", "#E63946", "High importance, underperforming → urgent action needed"),
            ("Strengths", "#2EC4B6", "High importance, strong performance → protect and promote"),
            ("Nice-to-Haves", "#A8DADC", "Low importance, high performance → over-invested"),
            ("Low Priority", "#8C9BB2", "Low importance, low performance → monitor only"),
        ]
        self._add_textbox(
            slide, "Quadrant Guide",
            Inches(9.3), Inches(1.4), Inches(3.6), Inches(0.35),
            font_size=10, bold=True, color=DARK_BLUE,
        )
        for i, (label, color, desc) in enumerate(quadrant_notes):
            y = Inches(1.85) + i * Inches(0.85)
            self._add_rect(slide, Inches(9.3), y, Inches(0.22), Inches(0.22), _hex_to_rgb(color))
            self._add_textbox(
                slide, label,
                Inches(9.62), y - Inches(0.02), Inches(3.2), Inches(0.28),
                font_size=9.5, bold=True, color=DARK_BLUE,
            )
            self._add_textbox(
                slide, desc,
                Inches(9.62), y + Inches(0.25), Inches(3.2), Inches(0.45),
                font_size=8.5, color=MID_GRAY, word_wrap=True, italic=True,
            )

    def _slide_driver_detail(self, predictor: str, insight: str, rank: int):
        slide = self._blank_slide()
        row = self.r.importance.ranking[self.r.importance.ranking["predictor"] == predictor].iloc[0]
        qrow = self.r.quadrants.quadrant_df[self.r.quadrants.quadrant_df["predictor"] == predictor].iloc[0]

        quadrant_color_map = {
            "Priority Fixes": _hex_to_rgb("#E63946"),
            "Strengths": TEAL,
            "Nice-to-Haves": _hex_to_rgb("#A8DADC"),
            "Low Priority": MID_GRAY,
        }
        q_color = quadrant_color_map.get(qrow["quadrant"], MID_GRAY)

        self._header_bar(
            slide,
            f"Driver #{rank}: {row['label']}",
            f"Quadrant: {qrow['quadrant']}  ·  Importance: {row['importance_pct']:.1f}%"
        )

        # Quadrant badge
        self._add_rect(slide, Inches(11.1), Inches(0.15), Inches(2.0), Inches(0.7), q_color)
        self._add_textbox(
            slide, qrow["quadrant"],
            Inches(11.1), Inches(0.15), Inches(2.0), Inches(0.7),
            font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )

        # Insight text
        self._add_textbox(
            slide, "Insight",
            Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.35),
            font_size=11, bold=True, color=DARK_BLUE,
        )
        self._add_rect(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.7), LIGHT_GRAY)
        self._add_textbox(
            slide, insight,
            Inches(0.6), Inches(1.75), Inches(12.0), Inches(1.5),
            font_size=12, color=DARK_BLUE, word_wrap=True,
        )

        # Stats cards
        rrow = self.r.regression.coefficients[self.r.regression.coefficients["predictor"] == predictor].iloc[0]
        crow = self.r.correlations.pearson[self.r.correlations.pearson["predictor"] == predictor].iloc[0]

        card_data = [
            ("Relative Importance", f"{row['importance_pct']:.1f}%"),
            ("Importance Rank", f"#{row['rank']} of {len(self.r.predictors)}"),
            ("Pearson r", f"{crow['r']:.3f}"),
            ("Std. Coefficient β", f"{rrow['std_coef']:.3f}"),
            ("Significant", "Yes" if rrow["significant"] else "No"),
            ("Performance", f"{qrow['performance_raw']:.2f}"),
        ]

        card_w = Inches(2.1)
        card_gap = Inches(0.1)
        start_x = Inches(0.4)
        for i, (label, value) in enumerate(card_data):
            x = start_x + i * (card_w + card_gap)
            self._add_rect(slide, x, Inches(3.55), card_w, Inches(1.05), DARK_BLUE)
            self._add_textbox(
                slide, value,
                x + Inches(0.05), Inches(3.6), card_w - Inches(0.1), Inches(0.5),
                font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
            )
            self._add_textbox(
                slide, label,
                x + Inches(0.05), Inches(4.12), card_w - Inches(0.1), Inches(0.35),
                font_size=8, color=MID_GRAY, align=PP_ALIGN.CENTER,
            )

        # Mini detail chart
        key = f"driver_{predictor}"
        if key in self.charts:
            self._add_image_bytes(
                slide, self.charts[key],
                Inches(0.4), Inches(4.8),
                width=Inches(12.5),
            )

    def _slide_recommendations(self, recs: list[str]):
        slide = self._blank_slide()
        self._header_bar(slide, "Recommendations", "Prioritized actions based on Key Driver Analysis")

        for i, rec in enumerate(recs[:6]):
            y = Inches(1.35) + i * Inches(0.92)
            # Number
            self._add_rect(slide, Inches(0.35), y + Inches(0.08), Inches(0.4), Inches(0.4), TEAL)
            self._add_textbox(
                slide, str(i + 1),
                Inches(0.35), y + Inches(0.08), Inches(0.4), Inches(0.4),
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            )
            # Recommendation text
            self._add_textbox(
                slide, rec,
                Inches(0.9), y, Inches(12.0), Inches(0.8),
                font_size=11, color=DARK_BLUE, word_wrap=True,
            )
            # Divider line
            if i < len(recs) - 1:
                self._add_rect(slide, Inches(0.35), y + Inches(0.88), Inches(12.6), Inches(0.01), LIGHT_GRAY)

    def _slide_appendix_regression(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Appendix: Regression Results", "OLS model coefficients and significance")

        coef_df = self.r.regression.coefficients

        # Column headers
        headers = ["Driver", "Coef.", "Std. β", "t-stat", "p-value", "Sig."]
        col_widths = [Inches(3.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.0)]
        col_starts = [Inches(0.35)]
        for w in col_widths[:-1]:
            col_starts.append(col_starts[-1] + w)

        header_y = Inches(1.3)
        self._add_rect(slide, Inches(0.3), header_y, Inches(12.7), Inches(0.38), DARK_BLUE)
        for header, x, w in zip(headers, col_starts, col_widths):
            self._add_textbox(
                slide, header,
                x, header_y + Inches(0.04), w, Inches(0.3),
                font_size=9, bold=True, color=WHITE,
            )

        for i, (_, row) in enumerate(coef_df.iterrows()):
            y = Inches(1.68) + i * Inches(0.42)
            bg = LIGHT_GRAY if i % 2 == 0 else WHITE
            self._add_rect(slide, Inches(0.3), y, Inches(12.7), Inches(0.4), bg)

            values = [
                row["predictor"].replace("_", " ").title(),
                f"{row['coef']:.4f}",
                f"{row['std_coef']:.4f}",
                f"{row['t_stat']:.3f}",
                f"{row['p_value']:.4f}",
                "✓" if row["significant"] else "—",
            ]
            for val, x, w in zip(values, col_starts, col_widths):
                color = TEAL if val == "✓" else DARK_BLUE
                self._add_textbox(
                    slide, val,
                    x, y + Inches(0.04), w, Inches(0.32),
                    font_size=9, color=color,
                )

        # Model summary
        self._add_textbox(
            slide,
            f"R² = {self.r.meta['r_squared']:.4f}  ·  Adj. R² = {self.r.meta['adj_r_squared']:.4f}  ·  "
            f"F = {self.r.regression.f_statistic:.2f} (p = {self.r.regression.f_p_value:.4f})",
            Inches(0.35), Inches(6.9), Inches(12.0), Inches(0.3),
            font_size=9, color=MID_GRAY, italic=True,
        )

    def _slide_appendix_correlations(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Appendix: Correlation Analysis", "Pearson and Spearman correlations with outcome")

        if "correlation" in self.charts:
            self._add_image_bytes(
                slide, self.charts["correlation"],
                Inches(0.5), Inches(1.25),
                width=Inches(12.3),
            )

    def _slide_appendix_model_fit(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Appendix: Model Fit", "Overall model performance metrics")

        if "model_fit" in self.charts:
            self._add_image_bytes(
                slide, self.charts["model_fit"],
                Inches(2.5), Inches(2.0),
                width=Inches(8.5),
            )

        self._add_textbox(
            slide,
            (
                f"The model accounts for {self.r.meta['r_squared']*100:.1f}% of variance in "
                f"{self.r.target.replace('_', ' ').title()} (Adj. R² = {self.r.meta['adj_r_squared']:.3f}). "
                f"F-statistic = {self.r.regression.f_statistic:.2f} "
                f"(p = {self.r.regression.f_p_value:.4f})."
            ),
            Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.6),
            font_size=11, color=DARK_BLUE, word_wrap=True, align=PP_ALIGN.CENTER,
        )

    # ── Build & save ──────────────────────────────────────────────────────────

    def build(self, exec_summary: str, recommendations: list[str], driver_insights: dict[str, str]) -> Presentation:
        """Assemble all slides in order."""
        logger.info("Building PowerPoint report...")

        self._slide_title()
        self._slide_exec_summary(exec_summary)
        self._slide_methodology()
        self._slide_importance_chart()
        self._slide_quadrant()

        # One slide per driver (top N)
        ranking = self.r.importance.ranking
        for _, row in ranking.iterrows():
            pred = row["predictor"]
            insight = driver_insights.get(pred, "")
            self._slide_driver_detail(pred, insight, int(row["rank"]))

        self._slide_recommendations(recommendations)
        self._slide_appendix_regression()
        self._slide_appendix_correlations()
        self._slide_appendix_model_fit()

        logger.info("Report complete: %d slides", len(self.prs.slides))
        return self.prs

    def save(self, path: str | Path) -> Path:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        logger.info("Saved report to %s", path)
        return path
